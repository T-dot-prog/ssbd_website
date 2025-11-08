from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Load existing certificate template
prs = Presentation("F:/Certificate.pptx")
slide = prs.slides[0]

# Match size and position from your screenshot
left = Inches(1.9)
top = Inches(3.13)
width = Inches(7.34)
height = Inches(0.85)

# Add textbox at that exact position
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame

# Add italic, centered name text
p = tf.add_paragraph()
p.text = "Md. Abdullah Inzum Adib"   # participant name
p.font.name = "Slight"               # or "Times New Roman"
p.font.italic = False
p.font.size = Pt(28)
p.alignment = PP_ALIGN.CENTER

# Save as a new file
prs.save("updated_certificate.pptx")


"""
Test image via request module
"""

import requests

IMAGE_ID = "1l9JDaLmuk5unOqlApHf-_xmWtqd47bTR"
url = f"https://drive.google.com/uc?export=download&id={IMAGE_ID}"

response = requests.get(url)

with open("image.jpg", "wb") as f:
    f.write(response.content)

print("Image downloaded successfully!")
